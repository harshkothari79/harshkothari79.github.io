import json
import os
import re
from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPT_PATH = os.path.join(ROOT, 'History', 'Timeline.pptx')
OUT_JSON = os.path.join(ROOT, 'History', 'timeline_extracted.json')

# Match four-digit years, ranges, 2-digit years with apostrophes, and years with +
YEAR_RE = re.compile(r"((?:18|19|20)\d{2}(?:\+|\s*[-–—]\s*(?:\d{2,4}))?|\b['’]?\d{2}\b)")


def iter_texts_from_shape(shape):
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame:
        for p in shape.text_frame.paragraphs:
            txt = (p.text or '').strip()
            if txt:
                yield txt
    if hasattr(shape, 'has_table') and shape.has_table and shape.table:
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                t = (cell.text or '').strip()
                if t:
                    yield t
    if hasattr(shape, 'shapes'):
        try:
            for shp in shape.shapes:
                yield from iter_texts_from_shape(shp)
        except Exception:
            pass


def extract_text_from_ppt(ppt_path: str):
    prs = Presentation(ppt_path)
    lines = []
    slides_data = []
    for slide in prs.slides:
        slide_lines = []
        title = None
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title = (slide.shapes.title.text or '').strip()
                if title:
                    slide_lines.append(title)
        except Exception:
            pass
        for shape in slide.shapes:
            for txt in iter_texts_from_shape(shape):
                slide_lines.append(txt)
        cleaned = []
        for t in slide_lines:
            t = t.replace('\uF0B7', '')
            tt = re.sub(r"\s+", " ", t).strip()
            if tt and (not cleaned or cleaned[-1] != tt):
                cleaned.append(tt)
        slides_data.append({'title': title, 'lines': cleaned})
        lines.extend(cleaned)
    return slides_data, lines


def split_history_and_timeline(slides_data, lines):
    events = []
    history_chunks = []
    for line in lines:
        m = YEAR_RE.search(line)
        if m:
            year = m.group(1)
            # Trim leading year and separators
            desc = re.sub(r"^\s*" + re.escape(year) + r"\s*[-:\u2013\u2014\u2022]*\s*", "", line).strip()
            events.append({'year': year, 'text': desc or line})
        else:
            if len(line) > 20:
                history_chunks.append(line)
    # Fallback: if no events detected, create one per slide from title or first line
    if not events:
        for i, s in enumerate(slides_data, start=1):
            title = (s.get('title') or '').strip()
            first_line = ''
            for ln in s.get('lines') or []:
                if ln and ln != title:
                    first_line = ln
                    break
            # detect a year anywhere in title/first line
            ysrc = title or first_line
            ym = YEAR_RE.search(ysrc or '')
            year = ym.group(1) if ym else f"Slide {i}"
            text = first_line or title or f"Timeline item {i}"
            events.append({'year': year, 'text': text})
    # History summary: prefer the longest substantial chunk, else slide 1 lines joined
    history_summary = max(history_chunks, key=len) if history_chunks else ''
    if not history_summary and slides_data:
        primary = []
        for ln in slides_data[0].get('lines') or []:
            if len(ln) > 20:
                primary.append(ln)
        history_summary = ' '.join(primary)[:600]
    return history_summary, events


def main():
    if not os.path.exists(PPT_PATH):
        print(json.dumps({'error': f'PPT not found at {PPT_PATH}'}))
        return
    slides, lines = extract_text_from_ppt(PPT_PATH)
    history_summary, events = split_history_and_timeline(slides, lines)
    data = {
        'ppt': PPT_PATH,
        'history_summary': history_summary,
        'events': events,
        'slides': slides,
    }
    with open(OUT_JSON, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(json.dumps({'ok': True, 'json': OUT_JSON, 'history_len': len(history_summary), 'events_count': len(events)}, ensure_ascii=False))

if __name__ == '__main__':
    main()